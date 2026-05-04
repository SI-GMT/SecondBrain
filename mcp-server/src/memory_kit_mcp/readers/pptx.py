"""PPTX reader — extract slides as Markdown via python-pptx.

Mirrors ``scripts/doc-readers/read_pptx.py``: each slide is rendered as
``## Slide N — {title}`` followed by text frames in document order, with
indentation reflecting paragraph level.
"""

from __future__ import annotations

from pathlib import Path

from . import DocReaderDependencyError

_PLACEHOLDER_SHAPE_TYPE = 14  # MSO_SHAPE_TYPE.PLACEHOLDER


def _shape_text_lines(shape) -> list[str]:
    if not shape.has_text_frame:
        return []
    lines: list[str] = []
    for para in shape.text_frame.paragraphs:
        text = "".join(run.text for run in para.runs).strip()
        if not text:
            continue
        prefix = "  " * max(0, para.level)
        if para.level > 0 or shape.shape_type == _PLACEHOLDER_SHAPE_TYPE:
            lines.append(f"{prefix}- {text}")
        else:
            lines.append(text)
    return lines


def _slide_title(slide) -> str:
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        title = slide.shapes.title.text_frame.text.strip()
        if title:
            return title
    return ""


def extract(path: Path) -> tuple[str, list[str]]:
    try:
        from pptx import Presentation
    except ImportError as exc:
        raise DocReaderDependencyError(
            "python-pptx is required to read .pptx files. "
            "Install via: pip install memory-kit-mcp[doc-readers]"
        ) from exc

    warnings: list[str] = []
    prs = Presentation(str(path))
    out: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        title = _slide_title(slide)
        header = f"## Slide {idx}" + (f" — {title}" if title else "")
        body_lines: list[str] = []
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            body_lines.extend(_shape_text_lines(shape))
        if not body_lines and not title:
            continue
        out.append(header)
        if body_lines:
            out.append("")
            out.extend(body_lines)
        out.append("")
    content = "\n".join(out).strip()
    if not content:
        warnings.append("no textual content extracted")
    return content, warnings
