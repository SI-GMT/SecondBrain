#!/usr/bin/env python3
# /// script
# requires-python = ">=3.10"
# dependencies = [
#   "python-pptx>=0.6.23",
# ]
# ///
"""
read_pptx.py — extract textual content from a .pptx file as Markdown.

Stdout: Markdown (each slide rendered as `## Slide N — {title}` followed by
text frames in document order).
Stderr: warnings and errors.
Exit codes:
  0  success
  1  invocation error
  2  empty extraction

Convention: invoked by core/procedures/mem-doc.md via `uv run`.
    uv run scripts/doc-readers/read_pptx.py {path}
"""

import argparse
import sys
from pathlib import Path

if hasattr(sys.stdout, "reconfigure"):
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
if hasattr(sys.stderr, "reconfigure"):
    sys.stderr.reconfigure(encoding="utf-8", errors="replace")

try:
    from pptx import Presentation
    from pptx.util import Pt  # noqa: F401
except ImportError:
    print("Error: python-pptx not available. Run via `uv run`.", file=sys.stderr)
    sys.exit(1)


def shape_text_lines(shape) -> list[str]:
    if not shape.has_text_frame:
        return []
    lines: list[str] = []
    for para in shape.text_frame.paragraphs:
        text = "".join(run.text for run in para.runs).strip()
        if not text:
            continue
        prefix = "  " * max(0, para.level)
        lines.append(f"{prefix}- {text}" if para.level > 0 or shape.shape_type == 14 else text)
    return lines


def slide_title(slide) -> str:
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        title = slide.shapes.title.text_frame.text.strip()
        if title:
            return title
    return ""


def extract(path: Path) -> str:
    prs = Presentation(str(path))
    out: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        title = slide_title(slide)
        header = f"## Slide {idx}" + (f" — {title}" if title else "")
        body_lines: list[str] = []
        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            body_lines.extend(shape_text_lines(shape))
        if not body_lines and not title:
            continue
        out.append(header)
        if body_lines:
            out.append("")
            out.extend(body_lines)
        out.append("")
    return "\n".join(out).strip()


def main() -> int:
    parser = argparse.ArgumentParser(description="Extract .pptx as Markdown")
    parser.add_argument("path", help="Path to .pptx file")
    args = parser.parse_args()

    path = Path(args.path)
    if not path.exists():
        print(f"Error: file not found: {path}", file=sys.stderr)
        return 1
    if not path.is_file():
        print(f"Error: not a file: {path}", file=sys.stderr)
        return 1

    try:
        content = extract(path)
    except Exception as exc:
        print(f"Error parsing pptx: {exc}", file=sys.stderr)
        return 1

    if not content:
        print("Warning: no textual content extracted", file=sys.stderr)
        return 2

    print(content)
    return 0


if __name__ == "__main__":
    sys.exit(main())
